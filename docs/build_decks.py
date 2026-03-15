"""Build two AKF presentation decks as .pptx files."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand colors ──
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xA0, 0xAE, 0xC0)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
SUBTLE = RGBColor(0x64, 0x74, 0x8B)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, left, top, width, height, items, font_size=16,
                     color=WHITE, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=ACCENT_BLUE, text_color=WHITE):
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    table = table_shape.table
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = data[i][j]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = "Calibri"
                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = text_color
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x1A, 0x25, 0x3C)
    return table_shape


# ═══════════════════════════════════════════════════════════════
# DECK 1: VC PITCH DECK
# ═══════════════════════════════════════════════════════════════

def build_vc_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 1, 1.5, 11, 1.2, "AKF", font_size=60, color=ACCENT_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 2.8, 11, 0.8, "Agent Knowledge Format", font_size=36, color=WHITE,
                 bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 2, 4.0, 9, 0.6, "The trust metadata standard for every file AI touches.",
                 font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 2, 5.5, 9, 0.5, "EXIF for AI", font_size=28, color=ACCENT_GREEN,
                 bold=True, alignment=PP_ALIGN.CENTER)

    # ── Slide 2: The Problem ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "The Problem", font_size=36, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, 0.8, 1.5, 11, 0.6, "AI output is amnesiac. Every file AI generates carries zero memory of how it was made.",
                 font_size=20, color=WHITE)

    items = [
        "No trust signal  --  Is this claim 98% confident or a hallucination?",
        "No provenance  --  Which model? Which source? Who verified it?",
        "No classification  --  Is this public, internal, or confidential?",
        "",
        "Result: Every downstream system re-processes everything from scratch.",
        "Multi-agent pipelines waste 40-70% of tokens re-verifying claims.",
        "Compliance teams cannot audit AI-generated content.",
        "No intelligent routing  --  everything gets the same expensive model.",
    ]
    add_bullet_slide(slide, 0.8, 2.5, 11, 4.5, items, font_size=18, color=LIGHT_GRAY)

    # ── Slide 3: The Solution ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "The Solution: AKF", font_size=36, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, 0.8, 1.4, 11, 0.6,
                 "Trust metadata that travels with the file. ~15 tokens of JSON.",
                 font_size=20, color=WHITE)

    data = [
        ["Question", "Today", "With AKF"],
        ["How confident?", "Unknown", "Per-claim score (0-1) with source tier"],
        ["Where from?", "Lost after generation", "Persistent provenance chain"],
        ["Who can see it?", "Manual org-level labels", "Embedded, inheritable, auditable"],
        ["AI or human?", "Binary guess", "Per-claim flag with model ID"],
        ["Still valid?", "No expiry", "Trust decay + TTL + freshness"],
    ]
    add_table(slide, 0.8, 2.3, 11.5, 4.0, 6, 3, data)

    # ── Slide 4: How It Works ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "How It Works", font_size=36, color=ACCENT_BLUE, bold=True)

    add_text_box(slide, 0.8, 1.5, 5.5, 0.5, "Compact (~15 tokens):", font_size=16, color=ACCENT_GREEN, bold=True)
    add_text_box(slide, 0.8, 2.0, 6, 1.2,
                 '{"v":"1.0","claims":[{"c":"Revenue $4.2B",\n "t":0.98,"src":"SEC 10-Q","tier":1}]}',
                 font_size=14, color=LIGHT_GRAY, font_name="Courier New")

    add_text_box(slide, 0.8, 3.5, 5.5, 0.5, "Embeds into everything:", font_size=16, color=ACCENT_GREEN, bold=True)
    items = [
        "DOCX / XLSX / PPTX  --  OOXML custom properties (visible in Word)",
        "PDF  --  metadata stream",
        "HTML  --  JSON-LD script tag",
        "Images  --  EXIF/XMP metadata",
        "JSON / Markdown / Code  --  inline metadata",
        "Everything else  --  sidecar .akf.json companion",
    ]
    add_bullet_slide(slide, 0.8, 4.1, 6, 3.0, items, font_size=14, color=LIGHT_GRAY)

    add_text_box(slide, 7.5, 1.5, 5, 0.5, "One-line API:", font_size=16, color=ACCENT_GREEN, bold=True)
    code_lines = [
        'import akf',
        '',
        '# Stamp trust metadata',
        'akf.stamp("Revenue $4.2B",',
        '  confidence=0.98, source="SEC 10-Q")',
        '',
        '# Embed into any file',
        'akf.embed("report.docx", claims=[...])',
        '',
        '# Audit for compliance',
        'akf.audit("report.akf",',
        '  regulation="eu_ai_act")',
    ]
    add_bullet_slide(slide, 7.5, 2.0, 5.2, 4.5, code_lines, font_size=13, color=LIGHT_GRAY,
                     font_name="Courier New")

    # ── Slide 5: Market Opportunity ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "Market Opportunity", font_size=36, color=ACCENT_BLUE, bold=True)

    items = [
        "AI infrastructure market: $100B+ by 2028",
        "Every AI-generated file needs trust metadata  --  billions of files/day",
        "EU AI Act (2026 enforcement) mandates transparency for AI content",
        "SOX, HIPAA, NIST AI RMF all require audit trails  --  AKF provides them",
        "",
        "Three monetization vectors:",
        "  1. Open-source SDK (adoption driver)  --  pip install akf",
        "  2. Enterprise API (trust-as-a-service)  --  hosted validation, routing, audit",
        "  3. Compliance platform  --  continuous governance for regulated industries",
    ]
    add_bullet_slide(slide, 0.8, 1.5, 11, 5.5, items, font_size=18, color=LIGHT_GRAY)

    # ── Slide 6: Competitive Landscape ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "Competitive Landscape", font_size=36, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, 0.8, 1.3, 11, 0.5, "Nobody solves claim-level trust for arbitrary file formats.",
                 font_size=18, color=WHITE, bold=True)

    data = [
        ["", "C2PA / Content Credentials", "Microsoft Purview", "Google SynthID", "AKF"],
        ["Scope", "Media only", "Microsoft 365", "Google AI only", "Any file format"],
        ["Trust level", "Binary (authentic/not)", "Org labels", "Binary (AI/not)", "Per-claim 0-1 score"],
        ["Provenance", "Signer chain", "None", "None", "Full hop chain"],
        ["Compliance", "No", "DLP only", "No", "EU AI Act, SOX, HIPAA, NIST"],
        ["Portable", "Media formats", "Vendor-locked", "Vendor-locked", "Open standard"],
        ["Cost to adopt", "PKI infrastructure", "E5 license", "Google Cloud", "pip install akf"],
    ]
    add_table(slide, 0.5, 2.0, 12.3, 4.5, 7, 5, data)

    # ── Slide 7: Value Proposition ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "Value Delivered", font_size=36, color=ACCENT_BLUE, bold=True)

    data = [
        ["Impact", "How", "Magnitude"],
        ["Token burn reduction", "Skip re-processing verified claims", "40-70% fewer tokens"],
        ["Compute savings", "Trust-aware routing to right-sized models", "50-60% compute reduction"],
        ["Compliance automation", "One-command audit for 6 regulations", "Weeks to minutes"],
        ["Multi-agent efficiency", "Trust signals survive agent handoffs", "Eliminate re-verification"],
        ["Data center energy", "Less redundant inference", "Proportional to token savings"],
    ]
    add_table(slide, 0.5, 1.5, 12.3, 4.5, 6, 3, data)

    # ── Slide 8: Traction & Progress ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "Traction", font_size=36, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, 0.8, 1.3, 11, 0.5, "Built in 15 days. Production-ready.",
                 font_size=20, color=WHITE, bold=True)

    items = [
        "28,500 lines of code  --  dual Python + TypeScript SDKs",
        "1,302 tests passing  --  comprehensive coverage",
        "32 CLI commands  --  full-featured developer tool",
        "9 file formats  --  native embed/extract (DOCX, PDF, XLSX, images...)",
        "4 framework integrations  --  LangChain, CrewAI, LlamaIndex, MCP Server",
        "4 extensions  --  VS Code, GitHub Action, Google Workspace, Office Add-in",
        "Published on PyPI (akf) and npm (akf-format)",
        "Live website: akf.dev",
        "JSON Schema: akf.dev/schema/v1.1",
    ]
    add_bullet_slide(slide, 0.8, 2.0, 11, 5.0, items, font_size=18, color=LIGHT_GRAY)

    # ── Slide 9: Business Model ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "Business Model", font_size=36, color=ACCENT_BLUE, bold=True)

    data = [
        ["Tier", "What", "Revenue"],
        ["Open Source (Free)", "SDK, CLI, file format spec, schema", "Adoption & ecosystem"],
        ["Pro API ($)", "Hosted validation, trust routing API, analytics", "Usage-based SaaS"],
        ["Enterprise ($$)", "Continuous compliance, custom integrations, SLA", "Annual contracts"],
        ["Standard License ($$$)", "OEM licensing for AI platforms embedding AKF", "Per-seat / per-API-call"],
    ]
    add_table(slide, 0.5, 1.5, 12.3, 4.2, 5, 3, data)

    add_text_box(slide, 0.8, 6.0, 11, 0.5,
                 "Playbook: Open standard drives adoption -> API monetizes scale -> Enterprise monetizes compliance.",
                 font_size=16, color=LIGHT_GRAY)

    # ── Slide 10: The Ask ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 1, 1.5, 11, 1.0, "The Ask", font_size=44, color=ACCENT_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    items = [
        "Seed round to fund:",
        "",
        "  1.  Developer advocacy & community (launch, content, integrations)",
        "  2.  Hosted API infrastructure (trust routing, validation-as-a-service)",
        "  3.  Enterprise pilot program (3-5 design partners in regulated industries)",
        "  4.  Team (2 engineers, 1 DevRel)",
    ]
    add_bullet_slide(slide, 2, 3.0, 9, 3.5, items, font_size=20, color=LIGHT_GRAY)

    add_text_box(slide, 2, 6.0, 9, 0.5, "akf.dev  |  pip install akf  |  npm install akf-format",
                 font_size=18, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

    prs.save("/Users/arun/akf/docs/AKF-VC-Pitch-Deck.pptx")
    print("VC deck saved: docs/AKF-VC-Pitch-Deck.pptx (10 slides)")


# ═══════════════════════════════════════════════════════════════
# DECK 2: TECHNICAL DECK (GitHub Repo)
# ═══════════════════════════════════════════════════════════════

def build_tech_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 1, 1.2, 11, 1.0, "AKF", font_size=60, color=ACCENT_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 2.5, 11, 0.8, "Technical Architecture Deep Dive",
                 font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 2, 3.8, 9, 0.5, "28.5K LOC  |  1,302 tests  |  32 CLI commands  |  Built in 15 days",
                 font_size=18, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 2, 5.0, 9, 0.5, "github.com/HMAKT99/AKF",
                 font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Slide 2: Repository Structure ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "Repository Structure", font_size=36, color=ACCENT_BLUE, bold=True)

    left_items = [
        "python/akf/          Python SDK (Pydantic v2)",
        "  models.py           10+ data models",
        "  cli.py              32 commands (Click)",
        "  stamp.py            One-line stamping API",
        "  universal.py        Format-agnostic embed/extract",
        "  trust.py            Trust computation engine",
        "  security.py         Security scoring (0-10)",
        "  report.py           Governance report renderer",
        "  formats/            Format-specific handlers",
        "  signing.py          Ed25519 cryptographic signing",
    ]
    add_bullet_slide(slide, 0.5, 1.5, 6, 5.5, left_items, font_size=13, color=LIGHT_GRAY,
                     font_name="Courier New")

    right_items = [
        "typescript/src/       TypeScript SDK (zero deps)",
        "  models.ts           Interfaces + type guards",
        "  normalize.ts        Compact <-> descriptive",
        "  trust.ts            Trust computation",
        "  compliance.ts       6 regulation checks",
        "",
        "packages/             Framework integrations",
        "  mcp-server-akf/     MCP server",
        "  langchain-akf/      LangChain handler",
        "  llama-index-akf/    LlamaIndex filter",
        "  crewai-akf/         CrewAI tool",
        "",
        "extensions/           Editor & CI",
        "  vscode/  github-action/  office-addin/",
    ]
    add_bullet_slide(slide, 6.8, 1.5, 6, 5.5, right_items, font_size=13, color=LIGHT_GRAY,
                     font_name="Courier New")

    # ── Slide 3: Data Model ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "Data Model (Pydantic v2 + TypeScript Interfaces)",
                 font_size=32, color=ACCENT_BLUE, bold=True)

    data = [
        ["Model", "Purpose", "Key Fields"],
        ["AKF (root)", "Top-level unit", "version, claims[], prov[], author, classification"],
        ["Claim", "Per-claim trust", "content, confidence, source, tier, ai_generated, verified"],
        ["ProvenanceHop", "Chain of custody", "hop, actor, action, timestamp, hash"],
        ["Origin", "Generation context", "type, model, provider, generation_params"],
        ["Review", "Human oversight", "reviewer, verdict, timestamp, comment"],
        ["Freshness", "Temporal validity", "ttl_days, decay_rate, assessed_at"],
        ["CostMetadata", "Token economics", "input_tokens, output_tokens, model, cost"],
        ["AgentProfile", "Agent identity", "agent_id, model, provider, capabilities"],
    ]
    add_table(slide, 0.5, 1.5, 12.3, 5.2, 9, 3, data)

    # ── Slide 4: Dual Format ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "Wire Format: Compact + Descriptive",
                 font_size=36, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, 0.8, 1.3, 11, 0.5,
                 "Same data, two representations. Lossless round-trip conversion.",
                 font_size=18, color=LIGHT_GRAY)

    add_text_box(slide, 0.8, 2.2, 5.5, 0.4, "Compact (~15 tokens, optimized for AI):",
                 font_size=16, color=ACCENT_GREEN, bold=True)
    compact = [
        '{',
        '  "v": "1.0",',
        '  "claims": [{',
        '    "c": "Revenue $4.2B",',
        '    "t": 0.98,',
        '    "src": "SEC 10-Q",',
        '    "tier": 1,',
        '    "ver": true',
        '  }]',
        '}',
    ]
    add_bullet_slide(slide, 0.8, 2.7, 5.5, 4.0, compact, font_size=14, color=LIGHT_GRAY,
                     font_name="Courier New")

    add_text_box(slide, 7.0, 2.2, 5.5, 0.4, "Descriptive (human-readable):",
                 font_size=16, color=ACCENT_GREEN, bold=True)
    descriptive = [
        '{',
        '  "version": "1.0",',
        '  "claims": [{',
        '    "content": "Revenue $4.2B",',
        '    "confidence": 0.98,',
        '    "source": "SEC 10-Q",',
        '    "authority_tier": 1,',
        '    "verified": true',
        '  }]',
        '}',
    ]
    add_bullet_slide(slide, 7.0, 2.7, 5.5, 4.0, descriptive, font_size=14, color=LIGHT_GRAY,
                     font_name="Courier New")

    # ── Slide 5: Trust Engine ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "Trust Computation Engine",
                 font_size=36, color=ACCENT_BLUE, bold=True)

    add_text_box(slide, 0.8, 1.5, 11, 0.5,
                 "effective_trust = confidence x authority_weight x temporal_decay x (1 + penalty)",
                 font_size=18, color=ACCENT_GREEN, bold=True, font_name="Courier New")

    data = [
        ["Tier", "Weight", "Example Sources", "Decision Thresholds"],
        ["1", "1.00", "SEC filings, official records", ">= 0.7  ->  ACCEPT"],
        ["2", "0.85", "Analyst reports, peer-reviewed", ">= 0.4  ->  LOW"],
        ["3", "0.70", "News, industry reports", "< 0.4   ->  REJECT"],
        ["4", "0.50", "Internal estimates, CRM data", ""],
        ["5", "0.30", "AI inference, extrapolations", ""],
    ]
    add_table(slide, 0.8, 2.5, 11.5, 3.5, 6, 4, data)

    items = [
        "Additional modifiers: origin_weight, grounding_bonus, review_bonus, decay over TTL",
        "Security scoring: 11-point checklist -> 0-10 score with letter grade (A-F)",
        "10 detection classes: hallucination risk, knowledge laundering, trust degradation, provenance gaps...",
    ]
    add_bullet_slide(slide, 0.8, 6.0, 11, 1.5, items, font_size=14, color=LIGHT_GRAY)

    # ── Slide 6: Universal Embed/Extract ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "Universal Embed / Extract",
                 font_size=36, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, 0.8, 1.3, 11, 0.5,
                 "One API for all formats. Metadata travels inside the file.",
                 font_size=18, color=LIGHT_GRAY)

    data = [
        ["Format", "Storage Mechanism", "Visible In"],
        ["DOCX / XLSX / PPTX", "docProps/custom.xml (OOXML standard)", "File > Properties > Custom"],
        ["PDF", "PDF metadata stream", "Document properties"],
        ["HTML", 'JSON-LD <script type="application/akf+json">', "Page source / DOM"],
        ["PNG / JPG", "EXIF UserComment / XMP", "Image metadata viewers"],
        ["JSON", 'Reserved "_akf" key', "Any JSON parser"],
        ["Markdown", "YAML frontmatter", "File header"],
        ["Email (.eml)", "X-AKF custom header", "Email headers"],
        ["Everything else", "Sidecar .akf.json companion", "Companion file"],
    ]
    add_table(slide, 0.4, 2.0, 12.5, 5.0, 9, 3, data)

    # ── Slide 7: CLI ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "CLI: 32 Commands",
                 font_size=36, color=ACCENT_BLUE, bold=True)

    data = [
        ["Category", "Commands", "Purpose"],
        ["Create & Stamp", "create, stamp, init, quickstart", "Generate trust metadata"],
        ["Inspect & Validate", "inspect, validate, trust, security, explain", "Analyze trust scores"],
        ["Embed & Extract", "embed, extract, read, scan, sidecar, convert", "Universal file format support"],
        ["Compliance", "audit, report, batch", "Regulatory checks (6 frameworks)"],
        ["Cryptography", "keygen, sign, verify, hash", "Ed25519 signing + integrity"],
        ["Git Integration", "stamp (git), log --trust", "Trust-annotated commits"],
        ["Advanced", "stream, consume, calibrate, freshness, doctor", "Production workflows"],
    ]
    add_table(slide, 0.4, 1.5, 12.5, 4.8, 8, 3, data)

    # ── Slide 8: Test Coverage ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "Test Coverage: 1,302 Tests",
                 font_size=36, color=ACCENT_BLUE, bold=True)

    data = [
        ["SDK", "Tests", "Files", "Runtime", "Status"],
        ["Python", "1,114", "44", "35.6s", "All passing (12 skipped - optional deps)"],
        ["TypeScript", "188", "12", "2.1s", "All passing"],
        ["Total", "1,302", "56", "37.7s", "100% green"],
    ]
    add_table(slide, 0.8, 1.5, 11.5, 2.5, 4, 5, data)

    items = [
        "Test categories: unit, integration, CLI, round-trip, format-specific, security, compliance",
        "Cross-format round-trip: embed -> extract verified for all 9 formats",
        "CLI command tests: every command tested with expected output",
        "Edge cases: XSS prevention, large metadata, special characters, legacy migration",
        "Vertical pipelines: create -> embed -> read -> audit -> report (end-to-end)",
    ]
    add_bullet_slide(slide, 0.8, 4.3, 11, 3.0, items, font_size=16, color=LIGHT_GRAY)

    # ── Slide 9: Integrations Ecosystem ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "Ecosystem: 12 Integrations",
                 font_size=36, color=ACCENT_BLUE, bold=True)

    data = [
        ["Type", "Integration", "What It Does"],
        ["AI Framework", "LangChain callback handler", "Auto-stamp trust on LLM calls"],
        ["AI Framework", "LlamaIndex node parser + trust filter", "Trust-aware RAG retrieval"],
        ["AI Framework", "CrewAI tool", "Trust-aware multi-agent tasks"],
        ["Protocol", "MCP Server (4 tools)", "create, validate, scan, trust via MCP"],
        ["Editor", "VS Code extension", "Syntax highlight, hover, validation"],
        ["CI/CD", "GitHub Action", "Fail builds on untrusted claims"],
        ["Office", "Google Workspace add-on", "Docs, Sheets, Slides integration"],
        ["Office", "Office Add-in", "Word, Excel, PowerPoint integration"],
        ["Agent", "8 skill files (.md)", "Discoverable by any AI agent"],
    ]
    add_table(slide, 0.3, 1.5, 12.7, 5.2, 10, 3, data)

    # ── Slide 10: Architecture Decisions ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "Architecture Decisions",
                 font_size=36, color=ACCENT_BLUE, bold=True)

    data = [
        ["Decision", "Choice", "Why"],
        ["Data models", "Pydantic v2", "Validation, serialization, JSON Schema gen"],
        ["CLI framework", "Click", "Composable, testable, auto-help"],
        ["TS runtime deps", "Zero", "No supply chain risk, tree-shakeable"],
        ["OOXML storage", "docProps/custom.xml", "Word/Excel native UI visibility"],
        ["Signing", "Ed25519 (PyNaCl)", "Fast, compact, modern (not RSA)"],
        ["Wire format", "Dual compact/descriptive", "15 tokens for AI, readable for humans"],
        ["Schema", "JSON Schema 2020-12", "Widest tooling support"],
        ["Report renderer", "Registry pattern", "Extensible (add formats via decorator)"],
    ]
    add_table(slide, 0.4, 1.5, 12.5, 5.2, 9, 3, data)

    # ── Slide 11: What's Next ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "What's Next", font_size=36, color=ACCENT_BLUE, bold=True)

    items = [
        "Trust-aware routing API  --  route queries to right-sized models based on claim trust",
        "Self-healing metadata  --  AI agents auto-refresh stale claims on cron",
        "Outcome tracking  --  compare predictions to actuals, auto-calibrate trust per source",
        "Hosted validation  --  akf.dev/api/validate endpoint for CI/CD",
        "Schema at stable URL  --  akf.dev/schema/v1.1.json",
        "Cross-language interop demos  --  Python creates, TypeScript consumes",
        "VS Code inline decorations  --  trust score gutter icons",
    ]
    add_bullet_slide(slide, 0.8, 1.5, 11, 5.5, items, font_size=18, color=LIGHT_GRAY)

    add_text_box(slide, 2, 6.2, 9, 0.5, "github.com/HMAKT99/AKF  |  akf.dev  |  pip install akf",
                 font_size=18, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

    prs.save("/Users/arun/akf/docs/AKF-Technical-Deep-Dive.pptx")
    print("Tech deck saved: docs/AKF-Technical-Deep-Dive.pptx (11 slides)")


if __name__ == "__main__":
    build_vc_deck()
    build_tech_deck()
