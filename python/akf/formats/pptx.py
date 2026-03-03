"""AKF v1.0 — PPTX format handler.

Embeds AKF metadata into PowerPoint PPTX files using the OOXML Custom XML
Part mechanism. Basic embed/extract operations require NO external
dependencies (uses stdlib zipfile only). Advanced operations like
auto_enrich and create optionally use python-pptx for slide-level
introspection.

Usage:
    from akf.formats.pptx import embed, extract, is_enriched, scan

    embed("deck.pptx", {"akf": "1.0", "claims": [...]})
    meta = extract("deck.pptx")
"""

from typing import List, Optional

from .base import AKFFormatHandler, ScanReport


class PPTXHandler(AKFFormatHandler):
    """Handler for Microsoft PowerPoint PPTX files."""

    FORMAT_NAME: str = "PPTX"
    EXTENSIONS: List[str] = [".pptx"]
    MODE: str = "embedded"
    MECHANISM: str = "OOXML Custom XML Part"
    DEPENDENCIES: List[str] = []  # Basic ops use zipfile only

    def embed(self, filepath: str, metadata: dict) -> None:
        """Embed AKF metadata into a PPTX file.

        Uses stdlib zipfile to add a custom XML part to the OOXML archive.
        No python-pptx dependency required.
        """
        from ._ooxml import embed_in_ooxml

        embed_in_ooxml(filepath, metadata)

    def extract(self, filepath: str) -> Optional[dict]:
        """Extract AKF metadata from a PPTX file.

        Returns None if no AKF metadata is embedded.
        """
        from ._ooxml import extract_from_ooxml

        return extract_from_ooxml(filepath)

    def is_enriched(self, filepath: str) -> bool:
        """Check if a PPTX file has AKF metadata embedded."""
        from ._ooxml import is_ooxml_enriched

        return is_ooxml_enriched(filepath)

    def auto_enrich(
        self,
        filepath: str,
        agent_id: str,
        default_tier: int = 3,
        classification: Optional[str] = None,
    ) -> None:
        """Auto-enrich a PPTX file with AKF metadata.

        If python-pptx is available, extracts text from slide shapes and
        creates per-slide claims. Otherwise, embeds basic metadata only.
        """
        meta = self._build_auto_metadata(filepath, agent_id, default_tier, classification)

        # Try to extract slide text for claims if python-pptx is available
        try:
            from pptx import Presentation

            prs = Presentation(filepath)
            claims = []
            for slide_idx, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                texts.append(text)
                if texts:
                    combined = " | ".join(texts)
                    if len(combined) > 10:
                        claims.append(
                            {
                                "location": "slide:{}".format(slide_idx + 1),
                                "c": combined[:200],
                                "t": 0.7,
                                "src": agent_id,
                                "ai": True,
                                "tier": default_tier,
                            }
                        )
            if claims:
                meta["claims"] = claims
                meta["ai_contribution"] = 1.0
        except ImportError:
            pass  # No python-pptx — embed basic metadata only

        self.embed(filepath, meta)


# ---------------------------------------------------------------------------
# Module-level convenience API
# ---------------------------------------------------------------------------

_handler = PPTXHandler()
embed = _handler.embed
extract = _handler.extract
is_enriched = _handler.is_enriched
scan = _handler.scan
auto_enrich = _handler.auto_enrich


def create(
    output: str,
    slides: list,
    classification: Optional[str] = None,
    author: Optional[str] = None,
) -> None:
    """Create a new PPTX with AKF metadata from structured slide content.

    Requires python-pptx. Install with: pip install akf[pptx]

    Args:
        output: Output file path.
        slides: List of slide dicts. Each dict should have:
            - title: Slide title text
            - content: (optional) Body text or list of bullet points
            - akf: (optional) Dict with claim metadata (t, src, ai, etc.)
        classification: Optional security classification label.
        author: Optional author/agent identifier.

    Example::

        create("deck.pptx", [
            {"title": "Q3 Results", "content": "Revenue $4.2B", "akf": {"t": 0.98}},
            {"title": "Outlook", "content": ["Growth expected", "New markets"]},
        ])
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise ImportError(
            "Creating PPTX requires python-pptx. Install with: pip install akf[pptx]"
        )

    import hashlib
    from datetime import datetime, timezone

    prs = Presentation()
    claims = []

    for i, slide_data in enumerate(slides):
        title_text = slide_data.get("title", "")
        content = slide_data.get("content", "")

        # Use a title+content layout (layout index 1 is typically Title and Content)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        if slide.shapes.title is not None:
            slide.shapes.title.text = title_text

        # Set content in the body placeholder
        body_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body_placeholder = shape
                break

        if body_placeholder is not None and content:
            if isinstance(content, list):
                tf = body_placeholder.text_frame
                tf.text = content[0] if content else ""
                for bullet in content[1:]:
                    p = tf.add_paragraph()
                    p.text = bullet
            else:
                body_placeholder.text_frame.text = str(content)

        # Build claim text
        claim_text = title_text
        if content:
            if isinstance(content, list):
                claim_text += " | " + " | ".join(content)
            else:
                claim_text += " | " + str(content)

        if "akf" in slide_data:
            claim = {
                "location": "slide:{}".format(i + 1),
                "c": claim_text[:200],
            }
            claim.update(slide_data["akf"])
            claims.append(claim)

    prs.save(output)

    # Compute hash of the saved file
    now = datetime.now(timezone.utc).isoformat()
    with open(output, "rb") as f:
        file_hash = "sha256:" + hashlib.sha256(f.read()).hexdigest()

    # Build metadata
    metadata = {
        "akf": "1.0",
        "generated_at": now,
        "overall_trust": (
            sum(c.get("t", 0.5) for c in claims) / len(claims) if claims else 0.5
        ),
        "ai_contribution": (
            sum(1 for c in claims if c.get("ai")) / len(claims) if claims else 0.0
        ),
        "claims": claims,
        "provenance": [
            {
                "actor": author or "unknown",
                "action": "created",
                "at": now,
                "hash": file_hash,
            }
        ],
        "integrity_hash": file_hash,
    }

    if classification is not None:
        metadata["classification"] = classification

    embed(output, metadata)
